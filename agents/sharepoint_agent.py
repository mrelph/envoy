"""SharePoint/OneDrive agent — search, browse, read, write files and manage lists."""

import asyncio
import os
import json
from agents.base import sharepoint

# Max chars to return from file reads to avoid blowing context windows
_MAX_TEXT = 80000


def _extract_docx(path: str) -> str:
    """Extract text from a .docx file."""
    from docx import Document
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path: str) -> str:
    """Extract text from a .pptx file."""
    from pptx import Presentation
    texts = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(t for t in texts if t.strip())


def _extract_text(path: str) -> str:
    """Extract readable text from a downloaded file."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".docx":
            return _extract_docx(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext in (".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm", ".log", ".py", ".js", ".ts"):
            with open(path, "r", errors="replace") as f:
                return f.read()
    except Exception as e:
        return f"Error extracting text from {path}: {e}"
    return f"Binary file downloaded to {path} — cannot extract text inline."


def _truncate(text: str) -> str:
    if len(text) <= _MAX_TEXT:
        return text
    return text[:_MAX_TEXT] + f"\n\n[... truncated — showing {_MAX_TEXT:,} of {len(text):,} chars]"


async def search(query: str, row_limit: int = 20) -> str:
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_search", {"query": query, "rowLimit": row_limit})
        raw = result.content[0].text if result.content else "No results."
        return _truncate(raw)


async def list_files(library: str = "Documents", folder: str = "", personal: bool = True,
                     site_url: str = "", top: int = 30) -> str:
    args = {"libraryName": library, "personal": personal, "top": top}
    if folder:
        args["folderPath"] = folder
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_list_files", args)
        return result.content[0].text if result.content else "No files found."


async def read_file(server_relative_url: str, inline: bool = True, personal: bool = True,
                    site_url: str = "") -> str:
    ext = os.path.splitext(server_relative_url)[1].lower()
    is_binary = ext in (".docx", ".pptx", ".xlsx", ".pdf", ".png", ".jpg", ".zip")

    args = {"serverRelativeUrl": server_relative_url, "savePath": "/tmp",
            "personal": personal, "inline": not is_binary}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_read_file", args)
        raw = result.content[0].text if result.content else "Could not read file."

    if is_binary:
        # MCP saves to /tmp — find the file and extract text
        filename = os.path.basename(server_relative_url)
        local_path = os.path.join("/tmp", filename)
        # Parse saved path from MCP response if available
        try:
            data = json.loads(raw)
            if isinstance(data, dict) and data.get("savedTo"):
                local_path = data["savedTo"]
        except (json.JSONDecodeError, TypeError):
            pass
        if os.path.exists(local_path):
            return _truncate(_extract_text(local_path))
        return f"File downloaded but could not locate at {local_path}. MCP response: {raw[:500]}"

    return _truncate(raw)


async def write_file(library: str, file_name: str, content: str = "", folder: str = "",
                     personal: bool = True, site_url: str = "") -> str:
    args = {"libraryName": library, "fileName": file_name, "personal": personal}
    if content:
        args["content"] = content
    if folder:
        args["folderPath"] = folder
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_write_file", args)
        return result.content[0].text if result.content else "Write failed."


async def list_libraries(personal: bool = True, site_url: str = "") -> str:
    args = {"personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_list_libraries", args)
        return result.content[0].text if result.content else "No libraries found."


async def list_sites() -> str:
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_list_sites", {})
        return result.content[0].text if result.content else "No sites found."


# --- Lists ---

async def list_lists(personal: bool = True, site_url: str = "") -> str:
    args = {"personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_list_lists", args)
        return result.content[0].text if result.content else "No lists found."


async def list_items(list_title: str, personal: bool = True, site_url: str = "",
                     filter_expr: str = "", top: int = 50) -> str:
    args = {"listTitle": list_title, "personal": personal, "top": top}
    if site_url:
        args["siteUrl"] = site_url
    if filter_expr:
        args["filter"] = filter_expr
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_list_items", args)
        return result.content[0].text if result.content else "No items found."


async def create_item(list_title: str, fields: dict, personal: bool = True, site_url: str = "") -> str:
    args = {"listTitle": list_title, "fields": fields, "personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_create_item", args)
        return result.content[0].text if result.content else "Create failed."


async def update_item(list_title: str, item_id: int, fields: dict, personal: bool = True,
                      site_url: str = "") -> str:
    args = {"listTitle": list_title, "itemId": item_id, "fields": fields, "personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_update_item", args)
        return result.content[0].text if result.content else "Update failed."


async def delete_item(list_title: str, item_id: int, personal: bool = True, site_url: str = "") -> str:
    args = {"listTitle": list_title, "itemId": item_id, "personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_delete_item", args)
        return result.content[0].text if result.content else "Delete failed."


async def create_list(title: str, description: str = "", personal: bool = True, site_url: str = "") -> str:
    args = {"title": title, "personal": personal}
    if description:
        args["description"] = description
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_create_list", args)
        return result.content[0].text if result.content else "Create list failed."


async def delete_list(list_title: str, personal: bool = True, site_url: str = "") -> str:
    args = {"listTitle": list_title, "personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_delete_list", args)
        return result.content[0].text if result.content else "Delete list failed."


# --- File management ---

async def upload_file(library: str, file_name: str, source_path: str, folder: str = "",
                      personal: bool = True, site_url: str = "") -> str:
    args = {"libraryName": library, "fileName": file_name, "sourcePath": source_path, "personal": personal}
    if folder:
        args["folderPath"] = folder
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_write_file", args)
        return result.content[0].text if result.content else "Upload failed."


async def delete_file(server_relative_url: str, personal: bool = True, site_url: str = "") -> str:
    args = {"serverRelativeUrl": server_relative_url, "personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_delete_file", args)
        return result.content[0].text if result.content else "Delete failed."


async def create_folder(library: str, folder_path: str, personal: bool = True, site_url: str = "") -> str:
    args = {"libraryName": library, "folderPath": folder_path, "personal": personal}
    if site_url:
        args["siteUrl"] = site_url
    async with sharepoint() as session:
        result = await session.call_tool("sharepoint_create_folder", args)
        return result.content[0].text if result.content else "Create folder failed."
