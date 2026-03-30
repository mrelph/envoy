"""Export agent — convert markdown content to Word (.docx) and PowerPoint (.pptx)."""

import os
import re
from datetime import datetime


def _output_path(filename: str) -> str:
    out_dir = os.path.expanduser("~/.envoy/exports")
    os.makedirs(out_dir, exist_ok=True)
    return os.path.join(out_dir, filename)


def _default_name(ext: str) -> str:
    ts = datetime.now().strftime("%Y%m%d-%H%M")
    return f"envoy-{ts}.{ext}"


def _parse_sections(markdown: str) -> list:
    """Split markdown into (level, title, body) tuples."""
    sections = []
    current_level, current_title, current_body = 0, "", []

    for line in markdown.splitlines():
        m = re.match(r'^(#{1,4})\s+(.*)', line)
        if m:
            if current_title or current_body:
                sections.append((current_level, current_title, "\n".join(current_body).strip()))
            current_level = len(m.group(1))
            current_title = m.group(2).strip()
            current_body = []
        else:
            current_body.append(line)

    if current_title or current_body:
        sections.append((current_level, current_title, "\n".join(current_body).strip()))
    return sections


def to_docx(markdown: str, filename: str = "") -> str:
    """Convert markdown to a Word document. Returns the file path."""
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    filename = filename or _default_name("docx")
    path = _output_path(filename)
    doc = Document()

    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)

    sections = _parse_sections(markdown)
    if not sections:
        doc.add_paragraph(markdown)
    else:
        for level, title, body in sections:
            if title:
                heading_level = min(level, 4)
                doc.add_heading(title, level=heading_level)
            if body:
                for line in body.splitlines():
                    line = line.strip()
                    if not line:
                        continue
                    # Bullet points
                    if re.match(r'^[-*•]\s+', line):
                        text = re.sub(r'^[-*•]\s+', '', line)
                        p = doc.add_paragraph(text, style='List Bullet')
                    # Numbered items
                    elif re.match(r'^\d+\.\s+', line):
                        text = re.sub(r'^\d+\.\s+', '', line)
                        p = doc.add_paragraph(text, style='List Number')
                    # Bold lines (likely sub-headers)
                    elif line.startswith('**') and line.endswith('**'):
                        p = doc.add_paragraph()
                        run = p.add_run(line.strip('* '))
                        run.bold = True
                    else:
                        # Strip inline markdown
                        clean = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
                        clean = re.sub(r'\*(.*?)\*', r'\1', clean)
                        clean = re.sub(r'`(.*?)`', r'\1', clean)
                        doc.add_paragraph(clean)

    doc.save(path)
    return path


def to_pptx(markdown: str, filename: str = "", title: str = "Envoy Report") -> str:
    """Convert markdown to a PowerPoint deck. Each ## section becomes a slide. Returns file path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    filename = filename or _default_name("pptx")
    path = _output_path(filename)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = datetime.now().strftime("%B %d, %Y")

    sections = _parse_sections(markdown)
    if not sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Report"
        slide.placeholders[1].text = markdown[:3000]
    else:
        for level, section_title, body in sections:
            if level > 2 and not section_title:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section_title or "Details"

            if body:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                first = True
                for line in body.splitlines():
                    line = line.strip()
                    if not line:
                        continue
                    # Strip markdown formatting
                    clean = re.sub(r'^[-*•]\s+', '', line)
                    clean = re.sub(r'^\d+\.\s+', '', clean)
                    clean = re.sub(r'\*\*(.*?)\*\*', r'\1', clean)
                    clean = re.sub(r'\*(.*?)\*', r'\1', clean)
                    clean = re.sub(r'`(.*?)`', r'\1', clean)

                    if first:
                        tf.paragraphs[0].text = clean
                        tf.paragraphs[0].font.size = Pt(18)
                        first = False
                    else:
                        p = tf.add_paragraph()
                        p.text = clean
                        p.font.size = Pt(18)
                        # Indent bullets
                        if re.match(r'^[-*•]\s+', line) or re.match(r'^\d+\.\s+', line):
                            p.level = 1

    prs.save(path)
    return path
